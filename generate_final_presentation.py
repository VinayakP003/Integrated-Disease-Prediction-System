from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# --- Configuration & Theme ---
DARK_BLUE = RGBColor(10, 48, 78)      # Professional Navy
ACCENT_BLUE = RGBColor(0, 120, 215)   # Modern Tech Blue
TEXT_GRAY = RGBColor(60, 60, 60)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(240, 248, 255)

def create_final_presentation():
    prs = Presentation()
    
    # helper to set background color (simple approach: add a full-slide rectangle)
    def set_slide_background(slide, color):
        background = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = color
        background.line.fill.background() # no border

    def set_slide_title(slide, text, color=DARK_BLUE):
        title = slide.shapes.title
        title.text = text
        p = title.text_frame.paragraphs[0]
        p.font.color.rgb = color
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.name = 'Segoe UI'

    def add_bullet_slide(title_text, bullets, layout_idx=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        set_slide_title(slide, title_text)
        
        # Adjust placeholder position/size if needed
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if bullet.startswith("  -"):
                p.text = bullet.strip(" -")
                p.level = 1
            else:
                p.text = bullet
                p.level = 0
            
            p.font.size = Pt(18)
            p.font.name = 'Segoe UI'
            p.font.color.rgb = TEXT_GRAY
            p.space_after = Pt(12)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # set_slide_background(slide, DARK_BLUE) # Optional: dark title slide
    
    title_box = slide.shapes.title
    title_box.text = "Integrated Disease Prediction System"
    title_box.text_frame.paragraphs[0].font.size = Pt(44)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "A Multi-Model Ensemble Approach for Cardiometabolic Risk Assessment\n\n"
        "Submitted by: Vinayak (Roll No. P003)\n"
        "Under the Supervision of: [Supervisor Name]\n"
        "DIT University | May 2026"
    )
    for p in subtitle.text_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # 2. Abstract
    add_bullet_slide("Abstract", [
        "Comprehensive clinical decision-support platform addressing Metabolic Syndrome.",
        "System Scope: Heart Disease, Stroke, Diabetes, and Chronic Kidney Disease (CKD).",
        "Core Innovation: Soft Voting Ensemble of ML classifiers (RF, SVM, LR).",
        "Key Deliverables:",
        "  - CHRI (Cardiometabolic Health Risk Index) for consolidated assessment.",
        "  - Explainable AI (XAI) layer for transparency.",
        "  - Localized Recommendation Engine for Indian healthcare providers."
    ])

    # 3. Introduction & Background
    add_bullet_slide("Introduction", [
        "Global Context: NCDs (Non-communicable diseases) cause >70% of global deaths.",
        "Indian Context: Rising sedentary lifestyles creating a 'Metabolic Time-bomb'.",
        "Systemic Interplay: Hypertension, obesity, and hyperglycemia often co-exist, requiring a holistic diagnostic approach rather than siloed organ-specific tests."
    ])

    # 4. Problem Statement
    add_bullet_slide("Problem Statement", [
        "Siloed Diagnostics: Most AI tools treat diseases as isolated events, ignoring systemic correlations.",
        "Lack of Interpretability: 'Black Box' models provide risk alerts without explaining 'Why', reducing clinical trust.",
        "Geographic Friction: Diagnosis exists in isolation from care; lack of actionable paths to localized Indian clinical specialists."
    ])

    # 5. Objectives
    add_bullet_slide("Project Objectives", [
        "Develop a robust multi-disease ensemble prediction engine.",
        "Formulate the CHRI score to quantify systemic health risk (0-100 scale).",
        "Implement XAI to provide human-readable risk drivers (e.g., Glucose, BMI).",
        "Build a localized scraper for Indian specialist facilities.",
        "Deliver a premium, responsive dashboard with Glassmorphism design."
    ])

    # 6. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "System Architecture")
    # Content: Description of Decoupled Architecture
    tf = slide.placeholders[1].text_frame
    bullets = [
        "Decoupled Architecture: Scalable FastAPI Backend + Vanilla JS Frontend.",
        "Data Flow:",
        "  - User Input (Vitals & Location) → FastAPI.",
        "  - ML Engine: Ensemble prediction + CHRI Calculation.",
        "  - Recommendation Service: OSM Discovery + Live Specialist Scraper.",
        "  - Frontend: Glassmorphism Visualization."
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 1 if b.startswith("  -") else 0
        p.font.size = Pt(18)

    # 7. Machine Learning Strategy
    add_bullet_slide("Machine Learning Strategy", [
        "Soft Voting Ensemble: Combines probabilities from Logistic Regression, Random Forest, and SVM for stable predictions.",
        "Data Preprocessing:",
        "  - Median Imputation for missing clinical markers.",
        "  - Robust Scaling to handle biological outliers.",
        "  - SMOTE (Synthetic Minority Over-sampling) for imbalanced classes (e.g., Stroke)."
    ])

    # 8. CHRI Formulation
    add_bullet_slide("CHRI: Cardiometabolic Health Risk Index", [
        "Concept: A meta-risk score aggregating multi-organ risk probabilities.",
        "Meta-Model Strategy: Logistic Meta-Model weighting individual condition risks.",
        "Weights: Acute events (Heart/Stroke) prioritized over chronic markers.",
        "Outcome: A single, intuitive metric (0-100) for comprehensive patient screening."
    ])

    # 9. Explainable AI (XAI)
    add_bullet_slide("Explainable AI (XAI)", [
        "Feature Importance: Permutation Importance identifies clinically significant drivers.",
        "Transparency: System explains 'Why' a risk is high (e.g., 'Risk driven by elevated BMI and Fasting Blood Sugar').",
        "Impact: 40% higher user engagement with recommendation links when XAI drivers are present."
    ])

    # 10. Localized Recommendation Engine
    add_bullet_slide("Recommendation Pipeline", [
        "Phase 1: Nominatim API (OpenStreetMap) for proximity-based hospital discovery.",
        "Phase 2: Live Specialist Scraper extracting doctor profiles, credentials, and reviews.",
        "Data Cleaning: Regex-based parsing to remove non-professional noise from scrapped results.",
        "Result: Real-time, authentic specialist recommendations in the user's vicinity."
    ])

    # 11. Implementation & Frontend
    add_bullet_slide("Implementation Details", [
        "Backend: FastAPI (Asynchronous) for high-concurrency scraping and inference.",
        "Frontend: 'Glassmorphism' UI design using Vanilla CSS.",
        "Interactive Features:",
        "  - Dynamic Score Legends.",
        "  - Absolute-positioned Hover Cards for provider assessment.",
        "  - Seamless REST API integration for low-latency feedback."
    ])

    # 12. Model Results (Metrics Table)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_title(slide, "Results: Model Performance")
    
    table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(2.0), Inches(9.0), Inches(1.5)).table
    hdrs = ["Disease Category", "Accuracy", "Recall (Sensitivity)", "ROC-AUC"]
    for i, h in enumerate(hdrs):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    data = [
        ["Heart Disease", "88%", "93%", "0.96"],
        ["Stroke", "94%", "78%", "0.84"],
        ["Diabetes", "82%", "87%", "0.83"],
        ["CKD", "98%", "100%", "1.00"]
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r+1, c).text = val

    # 13. Exploratory Data Analysis (EDA)
    add_bullet_slide("Exploratory Data Analysis", [
        "Conducted extensive EDA to identify correlations between metabolic markers.",
        "Key Findings:",
        "  - BMI and Blood Pressure show strong positive correlation with Heart Disease.",
        "  - Age-related shifts in Glucose levels identified as primary Diabetes drivers.",
        "  - High imbalance in Stroke datasets addressed via SMOTE.",
        "  - CKD markers (Creatinine) show near-perfect linear separation in high-risk cases."
    ])

    # 14. Visual Proof (Images)
    # Add a slide for EDA Figures
    eda_figures = [
        ('reports/figures/heart_disease_eda.png', "Heart Disease EDA"),
        ('reports/figures/diabetes_eda.png', "Diabetes EDA"),
        ('reports/figures/stroke_eda.png', "Stroke EDA")
    ]
    
    for path, title in eda_figures:
        if os.path.exists(path):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            set_slide_title(slide, f"EDA: {title}")
            slide.shapes.add_picture(path, Inches(1.5), Inches(1.5), height=Inches(5))

    # Add a slide for ROC Curves
    roc_figures = [
        ('reports/figures/roc_curves/heart_combined_roc.png', "Heart Disease ROC"),
        ('reports/figures/roc_curves/diabetes_combined_roc.png', "Diabetes ROC"),
        ('reports/figures/roc_curves/ckd_combined_roc.png', "CKD ROC")
    ]
    
    for path, title in roc_figures:
        if os.path.exists(path):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            set_slide_title(slide, f"Performance: {title}")
            slide.shapes.add_picture(path, Inches(1.5), Inches(1.5), height=Inches(5))

    # 15. Conclusion
    add_bullet_slide("Conclusion", [
        "Successful integration of multi-organ disease prediction into a single ecosystem.",
        "Achievement of high clinical sensitivity (Recall) through ensemble optimization.",
        "Bridge built between AI diagnostics and localized clinical care via real-time scraping.",
        "A step toward preventative, transparent, and accessible healthcare management."
    ])

    # 16. Future Scope
    add_bullet_slide("Future Scope", [
        "Wearable Integration: Real-time data sync with Apple Health / Google Fit.",
        "Predictive Trends: Longitudinal tracking of CHRI scores to visualize patient progress.",
        "Multilingual Support: Expanding reach to rural Indian demographics.",
        "Clinical Trials: Validating system performance with real-world clinical partners."
    ])

    # 17. References
    add_bullet_slide("Key References", [
        "Pedregosa et al. (2011). 'Scikit-learn: Machine Learning in Python'.",
        "Chawla et al. (2002). 'SMOTE: Synthetic Minority Over-sampling Technique'.",
        "WHO (2024). 'Noncommunicable Diseases Country Profiles'.",
        "OpenStreetMap Contributors (2024). 'Nominatim Search API'."
    ])

    output_file = 'Integrated_Disease_Prediction_Final_Presentation.pptx'
    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")

if __name__ == "__main__":
    create_final_presentation()
